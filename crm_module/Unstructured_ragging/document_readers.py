# document_readers.py
import os
import io
import base64
import pdfplumber
from pdf2image import convert_from_path
from docx import Document
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from PIL import Image
import pytesseract
from bs4 import BeautifulSoup
from config import TESSERACT_CMD, POPPLER_PATH

pytesseract.pytesseract.tesseract_cmd = TESSERACT_CMD

class DocumentReader:
    """Base class for readers"""
    def read(self, file_path):
        raise NotImplementedError

class TxtReader(DocumentReader):
    def read(self, file_path):
        with open(file_path, "r", encoding="utf-8") as f:
            lines = [l.strip() for l in f if l.strip()]
        return lines, [], []

class PdfReader(DocumentReader):
    def __init__(self):
        self.poppler_path = POPPLER_PATH

    def read(self, file_path):
        lines, tables, images = [], [], []

        # Extract text and tables
        try:
            with pdfplumber.open(file_path) as pdf:
                for page in pdf.pages:
                    text = page.extract_text()
                    if text:
                        lines.extend([l.strip() for l in text.split("\n") if l.strip()])
                    for table in page.extract_tables():
                        tables.append(table or [])
        except Exception as e:
            print(f"[PDF text/table error] {file_path}: {e}")

        # Extract images and OCR
        try:
            pages_img = convert_from_path(file_path, dpi=200, poppler_path=self.poppler_path)
            for img in pages_img:
                buf = io.BytesIO()
                img.save(buf, format="PNG")
                images.append(base64.b64encode(buf.getvalue()).decode("utf-8"))

                ocr_text = pytesseract.image_to_string(img)
                if ocr_text.strip():
                    lines.extend([l.strip() for l in ocr_text.split("\n") if l.strip()])
        except Exception as e:
            print(f"[PDF image/OCR error] {file_path}: {e}")

        return lines, tables, images

class DocxReader(DocumentReader):
    def read(self, file_path):
        lines, tables, images = [], [], []
        try:
            doc = Document(file_path)
            lines = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
            tables = [[cell.text for cell in row.cells] for table in doc.tables for row in table.rows]
            images = ["[docx image placeholder]" for _ in doc.inline_shapes]
        except Exception as e:
            print(f"[DOCX error] {file_path}: {e}")
        return lines, tables, images

class ExcelReader(DocumentReader):
    def read(self, file_path):
        import pandas as pd
        lines, tables, images = [], [], []
        try:
            df = pd.read_excel(file_path, sheet_name=None)
            for _, sheet in df.items():
                for row in sheet.itertuples(index=False):
                    row_text = " | ".join(str(cell) for cell in row if str(cell).strip())
                    if row_text:
                        lines.append(row_text)
        except Exception as e:
            print(f"[Excel error] {file_path}: {e}")
        return lines, tables, images

class PptxReader(DocumentReader):
    def read(self, file_path):
        lines, tables, images = [], [], []
        try:
            prs = Presentation(file_path)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if shape.has_text_frame and shape.text.strip():
                        lines.append(shape.text.strip())
                    if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                        tables.append([[cell.text for cell in row.cells] for row in shape.table.rows])
                    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                        buf = io.BytesIO()
                        buf.write(shape.image.blob)
                        images.append(base64.b64encode(buf.getvalue()).decode("utf-8"))
        except Exception as e:
            print(f"[PPTX error] {file_path}: {e}")
        return lines, tables, images

class HtmlReader(DocumentReader):
    def read(self, file_path):
        lines, tables, images = [], [], []
        try:
            with open(file_path, "r", encoding="utf-8") as f:
                soup = BeautifulSoup(f, "html.parser")
                lines = [t.strip() for t in soup.get_text().split("\n") if t.strip()]
                for table in soup.find_all("table"):
                    rows = [[td.get_text(strip=True) for td in tr.find_all(["td","th"])] for tr in table.find_all("tr")]
                    if rows:
                        tables.append(rows)
                for img in soup.find_all("img"):
                    src = img.get("src")
                    if src and src.startswith("data:image"):
                        images.append(src.split(",")[1])
        except Exception as e:
            print(f"[HTML error] {file_path}: {e}")
        return lines, tables, images

# Factory to get reader by extension
class ReaderFactory:
    READERS = {
        "txt": TxtReader(),
        "csv": TxtReader(),
        "pdf": PdfReader(),
        "docx": DocxReader(),
        "xls": ExcelReader(),
        "xlsx": ExcelReader(),
        "pptx": PptxReader(),
        "html": HtmlReader(),
    }

    @classmethod
    def get_reader(cls, ext):
        return cls.READERS.get(ext.lower())
